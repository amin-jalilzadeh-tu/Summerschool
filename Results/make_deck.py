"""Build the 8-slide case-study deck from the verified findings and the four figures.

Every number on these slides is reproduced by q1..q4_*.py in this folder. Nothing is typed
in by hand that a script does not also print.

    python3 make_deck.py   ->   CaseStudy2_StudioDelta.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "CaseStudy2_StudioDelta.pptx"

# --- palette, matching style.py so figures and slides read as one system ------------
INK = RGBColor(0x18, 0x18, 0x1B)
MUTED = RGBColor(0x52, 0x52, 0x5B)
FAINT = RGBColor(0xA1, 0xA1, 0xAA)
RULE = RGBColor(0xE4, 0xE4, 0xE7)
WASH = RGBColor(0xFA, 0xFA, 0xFA)
BLUE = RGBColor(0x25, 0x63, 0xEB)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x05, 0x96, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
W, H = 13.333, 7.5
ML, MR = 0.62, 0.62
CW = W - ML - MR

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]


# --- primitives ---------------------------------------------------------------------
def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, first=False, space_after=4, space_before=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    return p


def run(p, text, size=12, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.text_frame.word_wrap = True
    return s


def header(slide, kicker, title, sub=None, accent=BLUE):
    tf = box(slide, ML, 0.34, CW, 0.3)
    p = para(tf, first=True, space_after=0)
    run(p, kicker.upper(), size=11, color=accent, bold=True)

    tf = box(slide, ML, 0.63, CW, 0.52)
    p = para(tf, first=True, space_after=0)
    run(p, title, size=25, color=INK, bold=True)

    rect(slide, ML, 1.22, 1.05, 0.035, fill=accent)

    if sub:
        tf = box(slide, ML, 1.34, CW, 0.3)
        p = para(tf, first=True, space_after=0)
        run(p, sub, size=12.5, color=MUTED)


def footer(slide, n):
    tf = box(slide, ML, 7.16, CW - 0.5, 0.26)
    p = para(tf, first=True, space_after=0)
    run(p, "EURO PhD Summer School on MCDA/MCDM  ·  Case Study 2  ·  Studio Delta, Rotterdam",
        size=9, color=FAINT)
    tf = box(slide, W - MR - 0.5, 7.16, 0.5, 0.26)
    p = para(tf, first=True, space_after=0)
    p.alignment = PP_ALIGN.RIGHT
    run(p, str(n), size=9, color=FAINT)


def finding(slide, x, y, w, tag, tag_col, head, body, h=0.82):
    """A compact finding block: coloured tag, bold claim, one line of evidence."""
    rect(slide, x, y, 0.045, h, fill=tag_col)
    tf = box(slide, x + 0.16, y - 0.03, w - 0.16, h)
    p = para(tf, first=True, space_after=1)
    run(p, tag, size=9, color=tag_col, bold=True)
    p = para(tf, space_after=1)
    run(p, head, size=11.5, color=INK, bold=True)
    p = para(tf, space_after=0, line=1.0)
    run(p, body, size=9.5, color=MUTED)


#: Figures are 2.5:1. Sized so the bottom edge clears the findings strip at FIG_BOT.
FIG_TOP, FIG_W = 1.70, 11.2
FIG_X = (W - FIG_W) / 2
FIG_BOT = FIG_TOP + FIG_W / 2.5          # 6.18


def figure_slide(slide, png, y=FIG_TOP, width=FIG_W):
    slide.shapes.add_picture(str(HERE / png), Inches((W - width) / 2), Inches(y),
                             width=Inches(width))
    return y + width / 2.5


# =====================================================================================
# 1 · Title
# =====================================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, 0.13, fill=BLUE)

tf = box(s, ML, 1.55, CW, 0.34)
p = para(tf, first=True, space_after=0)
run(p, "EURO PhD SUMMER SCHOOL ON MCDA/MCDM  ·  TU DELFT  ·  CASE STUDY 2",
    size=12, color=BLUE, bold=True)

tf = box(s, ML, 2.05, CW - 1.6, 1.5)
p = para(tf, first=True, space_after=6, line=0.98)
run(p, "Designing a Solar-Efficient\nResidential Building in Rotterdam", size=40, color=INK, bold=True)

tf = box(s, ML, 3.62, CW - 2.2, 0.75)
p = para(tf, first=True, space_after=0, line=1.15)
run(p, "Which of thirty configurations should Studio Delta recommend — and which "
       "of them could any given method have recommended at all?", size=16, color=MUTED)

rect(s, ML, 4.62, 1.05, 0.035, fill=BLUE)

cards = [
    ("30 → 10", "survive the four\nstakeholder minimums", BLUE),
    ("3 of 10", "unreachable by any\npositive weighted sum", AMBER),
    ("+0.06", "and the feasible\nset is empty", RED),
    ("1 of 4", "requirements eliminates\nnothing at all", GREEN),
]
cx, cw, gap = ML, 2.72, 0.29
for big, small, col in cards:
    rect(s, cx, 4.95, cw, 1.28, fill=WASH, line=RULE)
    tf = box(s, cx + 0.22, 5.13, cw - 0.4, 0.42)
    p = para(tf, first=True, space_after=0)
    run(p, big, size=23, color=col, bold=True)
    tf = box(s, cx + 0.22, 5.58, cw - 0.4, 0.6)
    p = para(tf, first=True, space_after=0, line=1.05)
    run(p, small, size=11, color=MUTED)
    cx += cw + gap

tf = box(s, ML, 6.52, CW, 0.3)
p = para(tf, first=True, space_after=0)
run(p, "All figures reproducible:  ", size=10.5, color=FAINT)
run(p, "q1_design_space.py · q2_stakeholders.py · q3_thresholds.py · q4_criteria_audit.py",
    size=10.5, color=MUTED, font="Consolas")
footer(s, 1)

# =====================================================================================
# 2 · The problem, and what the case actually supplies
# =====================================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Framing", "The problem, and what the case actually supplies",
       "Meeting 3: the design space is already generated — our task is to analyse it and recommend")

# left column
rect(s, ML, 1.85, 6.12, 4.9, fill=WASH, line=RULE)
tf = box(s, ML + 0.3, 2.05, 5.5, 4.5)

p = para(tf, first=True, space_after=7)
run(p, "What we are given", size=14, color=INK, bold=True)
for label, val in [
    ("Alternatives", "30 configurations, explicitly enumerated (Appendix C)"),
    ("Criteria", "4, all normalised to [0, 1], all to be maximised"),
    ("Requirements", "PV ≥ 0.70 · daylight ≥ 0.70 · compactness ≥ 0.75 · FSI ≥ 0.80"),
    ("Problematique", "choice (α) — recommend one or more  ·  Kadziński p.6"),
]:
    p = para(tf, space_after=5, line=1.05)
    run(p, f"{label}   ", size=11, color=BLUE, bold=True)
    run(p, val, size=11.5, color=INK)

p = para(tf, space_before=10, space_after=4)
run(p, "The one fact that governs everything after it", size=14, color=INK, bold=True)
p = para(tf, space_after=6, line=1.1)
run(p, "The case supplies four minimum levels and ", size=12.5, color=INK)
run(p, "no weights.", size=12.5, color=RED, bold=True)
run(p, " That is the entire preference information in the document.", size=12.5, color=INK)
p = para(tf, space_after=0, line=1.1)
run(p, "So every weight vector we use is ours — an elicitation assumption, not data. "
       "It is named on every slide it appears on.", size=11, color=MUTED, italic=True)

# right column — methods
tf = box(s, ML + 6.5, 1.9, CW - 6.5, 0.34)
p = para(tf, first=True, space_after=0)
run(p, "Method choice, and why each earns its place", size=14, color=INK, bold=True)

methods = [
    ("Pareto dominance", "Köksalan I p.5–6 · Słowiński p.3",
     "The only preference-free comparison. Shows the space is fully conflicting.", BLUE),
    ("Conjunctive screen", "Köksalan I p.15–16 · Figueira I p.17",
     "Meeting 4 states minimums, and failing one disqualifies — non-compensatory.", BLUE),
    ("Supported / unsupported (LP)", "Köksalan I p.7, p.12–13",
     "Asks what no scoring method can: is this design reachable at all?", AMBER),
    ("Weight-space sampling", "Köksalan I p.14 · Kadziński p.14",
     "A distribution over winners instead of one arbitrary answer.", BLUE),
    ("Achievement scalarising fn.", "Köksalan II p.28–31 (Wierzbicki 1980)",
     "The only method here that reaches the unsupported points.", AMBER),
]
y = 2.34
for name, src, why, col in methods:
    rect(s, ML + 6.5, y, 0.038, 0.78, fill=col)
    tf = box(s, ML + 6.62, y - 0.03, CW - 6.7, 0.8)
    p = para(tf, first=True, space_after=1)
    run(p, name + "   ", size=12, color=INK, bold=True)
    run(p, src, size=9, color=FAINT)
    p = para(tf, space_after=0, line=1.03)
    run(p, why, size=10.5, color=MUTED)
    y += 0.80

tf = box(s, ML + 6.5, y + 0.06, CW - 6.5, 0.55)
p = para(tf, first=True, space_after=0, line=1.08)
run(p, "Rejected: ", size=10.5, color=RED, bold=True)
run(p, "a weighted sum alone (compensatory — it would buy past a failed minimum); "
       "ELECTRE Tri-nB and DRSA (right problematique for next time, but need parameters "
       "or labels nobody supplied).", size=10.5, color=MUTED)
footer(s, 2)

# =====================================================================================
# 3 · Q1
# =====================================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Question 1  ·  the feasible design space",
       "Which designs are even recommendable — before anyone argues about weights",
       "Appendix C (30 × 4) screened against the four Meeting-4 minimums")
figure_slide(s, "Q1_design_space.png")

fy, fw = FIG_BOT + 0.10, (CW - 0.6) / 3
finding(s, ML, fy, fw, "DOMINANCE DECIDES NOTHING", BLUE,
        "All 30 are non-dominated. The minimums cut 30 → 10.",
        "0 dominated pairs in 870 comparisons — the engineering team already stripped them. "
        "The screen does the work, not dominance.")
finding(s, ML + fw + 0.3, fy, fw, "REACHABILITY", AMBER,
        "C15, C25, C28 are unsupported in this set.",
        "No positive weight vector selects them — proved by LP, one per design. "
        "500 000 draws find them zero times; knife-edge C11 wins 158.")
finding(s, ML + 2 * (fw + 0.3), fy, fw, "THE METHOD DECIDES", AMBER,
        "A reference point reaches what no weight can.",
        "Wierzbicki ASF with q = y(C15) selects C15 outright. Which designs are "
        "recommendable at all depends on the scalarisation family.")
footer(s, 3)

# =====================================================================================
# 4 · Q2
# =====================================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Question 2  ·  three stakeholders",
       "What survives the weighting assumption — and what does not",
       "The case states four minimums and no weights, so every weight vector here is ours and is named")
figure_slide(s, "Q2_stakeholders.png")

fy, fw = FIG_BOT + 0.10, (CW - 0.6) / 3
finding(s, ML, fy, fw, "TWO ACTORS ARE STABLE", GREEN,
        "Client → C8. Engineering → C14.",
        "C8 is invariant for FSI weight 0.40–0.85, gap widening 0.004 → 0.016. "
        "C14 holds once compactness weight clears 5/11.")
finding(s, ML + fw + 0.3, fy, fw, "ONE IS NOT", RED,
        "The municipality gives three winners: C21, C1, C6.",
        "Five plausible weightings, margins 0.001–0.005 — at or below the ±0.005 rounding "
        "of a two-decimal table. The instability is in our model, not the actor.")
finding(s, ML + 2 * (fw + 0.3), fy, fw, "AGGREGATING ACTORS", BLUE,
        "Nash → C6.  Max-min → C11.",
        "Do not average the weight vectors — Vetschera p.33 flips the winner under an "
        "admissible rescaling. Nash is scale-invariant; max-min is not.")
footer(s, 4)

# =====================================================================================
# 5 · Q3
# =====================================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Question 3  ·  tightening and relaxing",
       "The requirements sit six hundredths from infeasibility",
       "Parametric ε-constraint analysis — Köksalan II p.11: “changing εᵢ systematically… "
       "explore different parts of the frontier”")
figure_slide(s, "Q3_thresholds.png")

fy, fw = FIG_BOT + 0.10, (CW - 0.6) / 3
finding(s, ML, fy, fw, "HEADROOM", RED,
        "+0.06 on all four and the feasible set is empty.",
        "−0.05 admits 17, +0.05 leaves 3. The set is far tighter on the upside: no room to "
        "negotiate up, considerable room down.")
finding(s, ML + fw + 0.3, fy, fw, "ONE REQUIREMENT IS INERT", RED,
        "Compactness ≥ 0.75 eliminates nothing.",
        "Delete it and the same ten designs survive — at every tightening level from +0.00 "
        "to +0.05. The other three are load-bearing by exactly one design each.")
finding(s, ML + 2 * (fw + 0.3), fy, fw, "A HUNDREDTH REWRITES THE GEOMETRY", AMBER,
        "Daylight 0.70 → 0.71 deletes C8, the client's pick.",
        "C15 replaces it — unsupported among the ten, but supported once C8 is gone. "
        "Supportedness is a property within a set, not a fixed label.")
footer(s, 5)

# =====================================================================================
# 6 · Q4a — the audit
# =====================================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Question 4  ·  auditing the criteria family",
       "The criteria family does not pass its own audit",
       "Cinelli p.42–43 (completeness · non-redundancy · preferential independence · conciseness), "
       "p.83–85 (correlations)")
figure_slide(s, "Q4_criteria_audit.png")

fy, fw = FIG_BOT + 0.10, (CW - 0.6) / 3
finding(s, ML, fy, fw, "NEAR-COLLINEARITY", AMBER,
        "All six pairs |r| ≥ 0.807.  PC1 = 92.4 %.",
        "Four criteria, one dominant dimension — and it does not soften on the ten we "
        "recommend from (90.5 %). A flag, not a verdict.")
finding(s, ML + fw + 0.3, fy, fw, "EXHAUSTIVENESS FAILS", RED,
        "Neighbour shadowing is named twice, measured never.",
        "Case pp. 2 and 4 raise sunlight “for the surrounding properties”. All four criteria "
        "measure only this building — two designs can tie and differ outside the plot.")
finding(s, ML + 2 * (fw + 0.3), fy, fw, "IT DID NOT REPRODUCE", BLUE,
        "Real geometry and weather: PC1 ≈ 0.51.",
        "We recomputed Appendix B on 3DBAG + PVGIS TMY. PV~FSI persists but weaker "
        "(+0.55 to +0.68); the PV~daylight antagonism does not appear at all.")
footer(s, 6)

# =====================================================================================
# 7 · Recommendation
# =====================================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Recommendation", "A set of four, with the conditions under which each wins",
       "Score gaps among the ten are 0.001–0.012, on a table published to two decimals — "
       "a single “winner” is not defensible")

recs = [
    ("C6", "COMPROMISE", BLUE,
     "Nash bargaining solution across the three actors (0.29713). Equal-weight optimum "
     "(0.8175). Wins L₁ and L₂ distance to the ideal, ties C11 at L∞. Municipality's pick "
     "under two of five weightings."),
    ("C11", "COMPROMISE", RGBColor(0x08, 0x91, 0xB2),
     "Max-min solution — the design no actor ranks worst. Second on equal weights (0.8125). "
     "Supported, but on a knife edge: 158 wins in 500 000 weight draws."),
    ("C8", "CLIENT EXTREME", RED,
     "Wins the largest share of the weight simplex (38.3 %) and is invariant across every "
     "FSI-heavy profile. Survives only because daylight = 0.70 exactly."),
    ("C14", "ENGINEERING EXTREME", GREEN,
     "Second largest share (30.2 %). Robust once compactness weight exceeds 5/11. "
     "Survives only because PV = 0.70 exactly."),
]
y = 1.92
for cid, role, col, why in recs:
    rect(s, ML, y, 8.55, 1.0, fill=WASH, line=RULE)
    rect(s, ML, y, 0.055, 1.0, fill=col)
    tf = box(s, ML + 0.28, y + 0.14, 1.15, 0.5)
    p = para(tf, first=True, space_after=0)
    run(p, cid, size=21, color=col, bold=True)
    tf = box(s, ML + 1.45, y + 0.13, 6.9, 0.8)
    p = para(tf, first=True, space_after=2)
    run(p, role, size=9, color=col, bold=True)
    p = para(tf, space_after=0, line=1.05)
    run(p, why, size=10.5, color=MUTED)
    y += 1.08

# right rail
rect(s, ML + 8.85, 1.92, CW - 8.85, 4.32, fill=WHITE, line=RULE)
tf = box(s, ML + 9.08, 2.1, CW - 9.3, 4.0)
p = para(tf, first=True, space_after=6)
run(p, "What we do not claim", size=13.5, color=INK, bold=True)
for t in [
    "Appendix C's normalisation is not published — the 0.70 cutoffs are numbers on an "
    "undisclosed scale.",
    "The thirty are a representative sample of the frontier by an undisclosed rule, not "
    "the frontier itself.",
    "Every weight vector is ours. BWM is proposed as the proper elicitation (2n−3 = 5 "
    "comparisons), not performed.",
    "Sampling cannot prove unsupportedness — the LP is the proof, the 500 000 draws "
    "corroborate.",
    "Setting q = y(C15) demonstrates reachability; it is not evidence that C15 is the "
    "right building.",
]:
    p = para(tf, space_after=6, line=1.06)
    run(p, "— ", size=10.5, color=AMBER, bold=True)
    run(p, t, size=10.5, color=MUTED)

p = para(tf, space_before=6, space_after=0, line=1.08)
run(p, "Choosing among these four is a political decision for the stakeholders, "
       "not a technical one for us.", size=11, color=INK, bold=True)
footer(s, 7)

# =====================================================================================
# 8 · For the next project
# =====================================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Question 4  ·  reuse", "For the next project: a protocol, not a spreadsheet",
       "The deliverable that transfers is how the method was chosen — and one change to the "
       "problematique")

cols = [
    ("1", "Derive criteria from objectives", BLUE,
     "Not from what the optimiser happens to compute. Stakeholder interviews → fundamental "
     "and means objectives → features → alternatives.\n\nCinelli p.86–91 — and it is an "
     "architecture case (Matassino 2024, Leiden Orangery)."),
    ("2", "Choose the method from features", BLUE,
     "Run the case through MCDA-MSS and say the sentence almost nobody writes:\n\n"
     "“We used X because our problem has features a, b, c.”\n\n"
     "Kadziński p.5–15 · 156 characteristics · p.18 rule-based recommendation."),
    ("3", "Change the problematique", AMBER,
     "This project is choice (α). The recurring question is sorting (β): is this massing "
     "acceptable, marginal, or unacceptable?\n\nSorting is absolute evaluation, so it is "
     "rank-reversal immune — next year's designs cannot reshuffle this year's verdicts.\n\n"
     "ELECTRE Tri-nB · DRSA (no weights at all)."),
]
cx, cw = ML, (CW - 0.56) / 3
for num, head, col, body in cols:
    rect(s, cx, 1.92, cw, 3.2, fill=WASH, line=RULE)
    tf = box(s, cx + 0.26, 2.1, cw - 0.5, 0.4)
    p = para(tf, first=True, space_after=0)
    run(p, num, size=22, color=col, bold=True)
    tf = box(s, cx + 0.26, 2.55, cw - 0.5, 0.42)
    p = para(tf, first=True, space_after=0)
    run(p, head, size=14, color=INK, bold=True)
    tf = box(s, cx + 0.26, 3.06, cw - 0.5, 1.95)
    p = para(tf, first=True, space_after=0, line=1.1)
    run(p, body, size=10.5, color=MUTED)
    cx += cw + 0.28

rect(s, ML, 5.32, CW, 0.035, fill=RULE)

tf = box(s, ML, 5.5, CW - 0.2, 1.45)
p = para(tf, first=True, space_after=0, line=1.16)
run(p, "The four criteria were all things the voxel optimiser could compute. ", size=15, color=INK)
run(p, "The criterion the brief names twice — effect on the neighbours — is the one it "
       "could not, and it is missing.", size=15, color=RED, bold=True)
run(p, " Meanwhile the four it did compute collapse onto one axis (PC1 = 92.4 %), and that "
       "collapse did not reproduce on real Rotterdam geometry and weather. And one of the "
       "four negotiated requirements never bound.", size=15, color=INK)
footer(s, 8)

prs.save(OUT)
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
